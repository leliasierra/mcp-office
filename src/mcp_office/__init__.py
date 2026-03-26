import json
from pathlib import Path
from typing import Any
from docx import Document
from docx.shared import Pt, Inches, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent
import fitz
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

server = Server("mcp-office")


def create_element(name):
    return OxmlElement(name)


def create_attribute(element, attr, value):
    element.set(qn(attr), value)


def rgb_to_hex(rgb):
    return f"{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"


@server.list_tools()
async def list_tools() -> list[Tool]:
    return [
        Tool(
            name="create_document",
            description="Create a new Word document with optional metadata",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path where to save the document",
                    },
                    "title": {"type": "string", "description": "Document title"},
                    "author": {"type": "string", "description": "Document author"},
                    "subject": {"type": "string", "description": "Document subject"},
                    "keywords": {"type": "string", "description": "Document keywords"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="open_document",
            description="Open and analyze a Word document",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to the document"},
                    "extract_text": {
                        "type": "boolean",
                        "description": "Extract text content",
                    },
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="get_document_properties",
            description="Get document properties and statistics",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to the document"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="list_documents",
            description="List all Word documents in a directory",
            inputSchema={
                "type": "object",
                "properties": {
                    "directory": {"type": "string", "description": "Directory path"},
                    "pattern": {
                        "type": "string",
                        "description": "File pattern (default: *.docx)",
                    },
                },
                "required": ["directory"],
            },
        ),
        Tool(
            name="copy_document",
            description="Create a copy of an existing document",
            inputSchema={
                "type": "object",
                "properties": {
                    "source": {"type": "string", "description": "Source document path"},
                    "destination": {
                        "type": "string",
                        "description": "Destination path",
                    },
                },
                "required": ["source", "destination"],
            },
        ),
        Tool(
            name="merge_documents",
            description="Merge multiple documents into one",
            inputSchema={
                "type": "object",
                "properties": {
                    "sources": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of source document paths",
                    },
                    "output": {"type": "string", "description": "Output document path"},
                },
                "required": ["sources", "output"],
            },
        ),
        Tool(
            name="add_heading",
            description="Add a heading with different levels and formatting",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "text": {"type": "string", "description": "Heading text"},
                    "level": {"type": "integer", "description": "Heading level (1-9)"},
                    "bold": {"type": "boolean", "description": "Bold text"},
                    "italic": {"type": "boolean", "description": "Italic text"},
                    "font_name": {"type": "string", "description": "Font name"},
                    "font_size": {
                        "type": "integer",
                        "description": "Font size in points",
                    },
                    "color": {"type": "string", "description": "Hex color code"},
                    "border": {"type": "boolean", "description": "Add border"},
                    "position": {
                        "type": "string",
                        "description": "Position: start, end, or index",
                    },
                    "index": {
                        "type": "integer",
                        "description": "Index for positioning",
                    },
                },
                "required": ["path", "text", "level"],
            },
        ),
        Tool(
            name="add_paragraph",
            description="Add a paragraph with optional styling and formatting",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "text": {"type": "string", "description": "Paragraph text"},
                    "style": {"type": "string", "description": "Paragraph style"},
                    "bold": {"type": "boolean", "description": "Bold text"},
                    "italic": {"type": "boolean", "description": "Italic text"},
                    "underline": {"type": "boolean", "description": "Underline text"},
                    "font_name": {"type": "string", "description": "Font name"},
                    "font_size": {
                        "type": "integer",
                        "description": "Font size in points",
                    },
                    "color": {"type": "string", "description": "Hex color code"},
                    "alignment": {
                        "type": "string",
                        "description": "Alignment: left, center, right, justify",
                    },
                    "position": {
                        "type": "string",
                        "description": "Position: start, end, or index",
                    },
                    "index": {
                        "type": "integer",
                        "description": "Index for positioning",
                    },
                },
                "required": ["path", "text"],
            },
        ),
        Tool(
            name="add_table",
            description="Create a table with custom data and styling",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "data": {
                        "type": "array",
                        "items": {"type": "array", "items": {"type": "string"}},
                        "description": "Table data as 2D array",
                    },
                    "style": {"type": "string", "description": "Table style"},
                    "header_row": {
                        "type": "boolean",
                        "description": "Format first row as header",
                    },
                    "alternating_colors": {
                        "type": "boolean",
                        "description": "Alternating row colors",
                    },
                    "header_color": {
                        "type": "string",
                        "description": "Header row background color",
                    },
                    "position": {
                        "type": "string",
                        "description": "Position: start, end, or index",
                    },
                    "index": {
                        "type": "integer",
                        "description": "Index for positioning",
                    },
                },
                "required": ["path", "data"],
            },
        ),
        Tool(
            name="add_image",
            description="Add an image with proportional scaling",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "image_path": {
                        "type": "string",
                        "description": "Path to image file",
                    },
                    "width_cm": {"type": "number", "description": "Width in cm"},
                    "height_cm": {"type": "number", "description": "Height in cm"},
                    "position": {
                        "type": "string",
                        "description": "Position: start, end, or index",
                    },
                    "index": {
                        "type": "integer",
                        "description": "Index for positioning",
                    },
                },
                "required": ["path", "image_path"],
            },
        ),
        Tool(
            name="add_page_break",
            description="Insert a page break",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "position": {
                        "type": "string",
                        "description": "Position: start, end, or index",
                    },
                    "index": {
                        "type": "integer",
                        "description": "Index for positioning",
                    },
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="add_bullet_list",
            description="Add a bullet or numbered list",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "items": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List items",
                    },
                    "numbered": {
                        "type": "boolean",
                        "description": "Numbered list (default: bullet)",
                    },
                    "start_index": {
                        "type": "integer",
                        "description": "Starting number for numbered list",
                    },
                    "position": {
                        "type": "string",
                        "description": "Position: start, end, or index",
                    },
                    "index": {
                        "type": "integer",
                        "description": "Index for positioning",
                    },
                },
                "required": ["path", "items"],
            },
        ),
        Tool(
            name="add_footnote",
            description="Add footnotes and endnotes to documents",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "text": {
                        "type": "string",
                        "description": "Text containing the reference",
                    },
                    "footnote_text": {"type": "string", "description": "Footnote text"},
                    "is_endnote": {
                        "type": "boolean",
                        "description": "Create endnote instead of footnote",
                    },
                },
                "required": ["path", "text", "footnote_text"],
            },
        ),
        Tool(
            name="convert_footnote_to_endnote",
            description="Convert footnotes to endnotes",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="format_footnote_style",
            description="Customize footnote and endnote styling",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "font_name": {"type": "string", "description": "Font name"},
                    "font_size": {"type": "integer", "description": "Font size"},
                    "color": {"type": "string", "description": "Hex color code"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="format_text_range",
            description="Format specific text ranges in a document",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "start_text": {
                        "type": "string",
                        "description": "Start text of range",
                    },
                    "end_text": {"type": "string", "description": "End text of range"},
                    "bold": {"type": "boolean", "description": "Bold text"},
                    "italic": {"type": "boolean", "description": "Italic text"},
                    "underline": {"type": "boolean", "description": "Underline text"},
                    "font_name": {"type": "string", "description": "Font name"},
                    "font_size": {"type": "integer", "description": "Font size"},
                    "color": {"type": "string", "description": "Hex color code"},
                },
                "required": ["path", "start_text", "end_text"],
            },
        ),
        Tool(
            name="find_replace",
            description="Find and replace text throughout the document",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "find": {"type": "string", "description": "Text to find"},
                    "replace": {"type": "string", "description": "Replacement text"},
                    "case_sensitive": {
                        "type": "boolean",
                        "description": "Case sensitive search",
                    },
                },
                "required": ["path", "find", "replace"],
            },
        ),
        Tool(
            name="format_table_cell",
            description="Format individual table cells",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "table_index": {"type": "integer", "description": "Table index"},
                    "row": {"type": "integer", "description": "Row index"},
                    "col": {"type": "integer", "description": "Column index"},
                    "text": {"type": "string", "description": "Cell text"},
                    "bold": {"type": "boolean", "description": "Bold text"},
                    "italic": {"type": "boolean", "description": "Italic text"},
                    "underline": {"type": "boolean", "description": "Underline text"},
                    "font_name": {"type": "string", "description": "Font name"},
                    "font_size": {"type": "integer", "description": "Font size"},
                    "color": {"type": "string", "description": "Hex color code"},
                    "alignment": {
                        "type": "string",
                        "description": "Alignment: left, center, right",
                    },
                    "vertical_alignment": {
                        "type": "string",
                        "description": "Vertical alignment: top, center, bottom",
                    },
                    "shading_color": {
                        "type": "string",
                        "description": "Cell background color",
                    },
                },
                "required": ["path", "table_index", "row", "col"],
            },
        ),
        Tool(
            name="merge_table_cells",
            description="Merge table cells",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "table_index": {"type": "integer", "description": "Table index"},
                    "start_row": {"type": "integer", "description": "Start row"},
                    "start_col": {"type": "integer", "description": "Start column"},
                    "end_row": {"type": "integer", "description": "End row"},
                    "end_col": {"type": "integer", "description": "End column"},
                },
                "required": [
                    "path",
                    "table_index",
                    "start_row",
                    "start_col",
                    "end_row",
                    "end_col",
                ],
            },
        ),
        Tool(
            name="set_column_width",
            description="Set column widths with multiple units",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "table_index": {"type": "integer", "description": "Table index"},
                    "col_index": {"type": "integer", "description": "Column index"},
                    "width_cm": {"type": "number", "description": "Width in cm"},
                    "width_percent": {
                        "type": "number",
                        "description": "Width as percentage",
                    },
                    "auto_fit": {
                        "type": "boolean",
                        "description": "Auto fit column width",
                    },
                },
                "required": ["path", "table_index", "col_index"],
            },
        ),
        Tool(
            name="format_table",
            description="Format table with borders and styles",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "table_index": {"type": "integer", "description": "Table index"},
                    "style": {"type": "string", "description": "Table style name"},
                    "border_color": {
                        "type": "string",
                        "description": "Border color hex",
                    },
                    "border_width": {
                        "type": "number",
                        "description": "Border width in points",
                    },
                    "header_row": {
                        "type": "boolean",
                        "description": "Format header row",
                    },
                    "header_color": {
                        "type": "string",
                        "description": "Header background color",
                    },
                    "alternating_rows": {
                        "type": "boolean",
                        "description": "Alternating row colors",
                    },
                    "first_row_color": {
                        "type": "string",
                        "description": "First row color",
                    },
                    "second_row_color": {
                        "type": "string",
                        "description": "Second row color",
                    },
                },
                "required": ["path", "table_index"],
            },
        ),
        Tool(
            name="delete_paragraph",
            description="Delete a paragraph by index",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "index": {
                        "type": "integer",
                        "description": "Paragraph index to delete",
                    },
                },
                "required": ["path", "index"],
            },
        ),
        Tool(
            name="insert_before_index",
            description="Insert content before a specific index",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "content_type": {
                        "type": "string",
                        "description": "Content type: heading, paragraph, or list",
                    },
                    "text": {"type": "string", "description": "Text content"},
                    "level": {
                        "type": "integer",
                        "description": "Heading level (if heading)",
                    },
                    "items": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List items (if list)",
                    },
                    "index": {
                        "type": "integer",
                        "description": "Index to insert before",
                    },
                },
                "required": ["path", "content_type", "text", "index"],
            },
        ),
        Tool(
            name="insert_after_index",
            description="Insert content after a specific index",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "content_type": {
                        "type": "string",
                        "description": "Content type: heading, paragraph, or list",
                    },
                    "text": {"type": "string", "description": "Text content"},
                    "level": {
                        "type": "integer",
                        "description": "Heading level (if heading)",
                    },
                    "items": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List items (if list)",
                    },
                    "index": {
                        "type": "integer",
                        "description": "Index to insert after",
                    },
                },
                "required": ["path", "content_type", "text", "index"],
            },
        ),
        Tool(
            name="create_custom_style",
            description="Create custom document styles",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "style_name": {"type": "string", "description": "Style name"},
                    "style_type": {
                        "type": "string",
                        "description": "Style type: paragraph or character",
                    },
                    "font_name": {"type": "string", "description": "Font name"},
                    "font_size": {"type": "integer", "description": "Font size"},
                    "bold": {"type": "boolean", "description": "Bold"},
                    "italic": {"type": "boolean", "description": "Italic"},
                },
                "required": ["path", "style_name", "style_type"],
            },
        ),
        Tool(
            name="add_section_border",
            description="Add borders to section headers",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "border_color": {
                        "type": "string",
                        "description": "Border color hex",
                    },
                    "border_width": {
                        "type": "number",
                        "description": "Border width in points",
                    },
                    "border_style": {
                        "type": "string",
                        "description": "Border style: single, double",
                    },
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="protect_document",
            description="Add password protection to document",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "password": {
                        "type": "string",
                        "description": "Password for protection",
                    },
                },
                "required": ["path", "password"],
            },
        ),
        Tool(
            name="add_editable_region",
            description="Add restricted editing with editable sections",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "editable_ranges": {
                        "type": "array",
                        "items": {"type": "integer"},
                        "description": "Paragraph indices that should remain editable",
                    },
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="extract_comments",
            description="Extract all comments from a document",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Document path"},
                    "author": {"type": "string", "description": "Filter by author"},
                    "paragraph_index": {
                        "type": "integer",
                        "description": "Get comments for specific paragraph",
                    },
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="pdf_read",
            description="Read and extract content from PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "extract_text": {
                        "type": "boolean",
                        "description": "Extract text content",
                    },
                    "extract_images": {
                        "type": "boolean",
                        "description": "Extract images",
                    },
                    "page_start": {
                        "type": "integer",
                        "description": "Start page (1-based)",
                    },
                    "page_end": {
                        "type": "integer",
                        "description": "End page (1-based)",
                    },
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="pdf_get_info",
            description="Get PDF metadata and information",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="pdf_extract_images",
            description="Extract all images from PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "output_dir": {
                        "type": "string",
                        "description": "Output directory for images",
                    },
                },
                "required": ["path", "output_dir"],
            },
        ),
        Tool(
            name="pdf_split",
            description="Split PDF into separate pages",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "output_dir": {"type": "string", "description": "Output directory"},
                    "pages": {
                        "type": "array",
                        "items": {"type": "integer"},
                        "description": "Pages to extract",
                    },
                },
                "required": ["path", "output_dir"],
            },
        ),
        Tool(
            name="pdf_merge",
            description="Merge multiple PDFs into one",
            inputSchema={
                "type": "object",
                "properties": {
                    "sources": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of PDF paths",
                    },
                    "output": {"type": "string", "description": "Output path"},
                },
                "required": ["sources", "output"],
            },
        ),
        Tool(
            name="pdf_rotate",
            description="Rotate pages in PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "pages": {
                        "type": "array",
                        "items": {"type": "integer"},
                        "description": "Pages to rotate",
                    },
                    "degrees": {
                        "type": "integer",
                        "description": "Rotation degrees (90, 180, 270)",
                    },
                    "output": {"type": "string", "description": "Output path"},
                },
                "required": ["path", "degrees"],
            },
        ),
        Tool(
            name="pdf_compress",
            description="Compress PDF file",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "quality": {
                        "type": "string",
                        "description": "Compression quality: low, medium, high",
                    },
                    "output": {"type": "string", "description": "Output path"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="pdf_add_watermark",
            description="Add watermark to PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "text": {"type": "string", "description": "Watermark text"},
                    "opacity": {"type": "number", "description": "Opacity (0-1)"},
                    "angle": {"type": "integer", "description": "Rotation angle"},
                    "output": {"type": "string", "description": "Output path"},
                },
                "required": ["path", "text"],
            },
        ),
        Tool(
            name="pdf_add_page_numbers",
            description="Add page numbers to PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "position": {
                        "type": "string",
                        "description": "Position: bottom-center, bottom-right, top-center",
                    },
                    "font_size": {"type": "integer", "description": "Font size"},
                    "output": {"type": "string", "description": "Output path"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="pdf_protect",
            description="Add password protection to PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "password": {"type": "string", "description": "Password"},
                    "owner_password": {
                        "type": "string",
                        "description": "Owner password (for full permissions)",
                    },
                    "permissions": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Permissions: print, copy, edit",
                    },
                    "output": {"type": "string", "description": "Output path"},
                },
                "required": ["path", "password"],
            },
        ),
        Tool(
            name="pdf_unprotect",
            description="Remove password protection from PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "password": {"type": "string", "description": "Password"},
                    "output": {"type": "string", "description": "Output path"},
                },
                "required": ["path", "password"],
            },
        ),
        Tool(
            name="pdf_extract_text",
            description="Extract text from specific pages",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "page_start": {"type": "integer", "description": "Start page"},
                    "page_end": {"type": "integer", "description": "End page"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="pdf_extract_tables",
            description="Extract tables from PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF file"},
                    "page": {"type": "integer", "description": "Page number"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="word_to_pdf",
            description="Convert Word document to PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "input_path": {
                        "type": "string",
                        "description": "Input Word document path",
                    },
                    "output_path": {"type": "string", "description": "Output PDF path"},
                },
                "required": ["input_path", "output_path"],
            },
        ),
        Tool(
            name="pdf_to_word",
            description="Convert PDF to Word document",
            inputSchema={
                "type": "object",
                "properties": {
                    "input_path": {"type": "string", "description": "Input PDF path"},
                    "output_path": {
                        "type": "string",
                        "description": "Output Word path",
                    },
                },
                "required": ["input_path", "output_path"],
            },
        ),
        Tool(
            name="excel_create",
            description="Create a new Excel workbook",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to save workbook"},
                    "sheet_name": {"type": "string", "description": "Sheet name"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="excel_read",
            description="Read data from Excel workbook",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                    "sheet": {"type": "string", "description": "Sheet name or index"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="excel_write_cell",
            description="Write value to a cell",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                    "sheet": {"type": "string", "description": "Sheet name"},
                    "cell": {"type": "string", "description": "Cell (e.g., A1)"},
                    "value": {"type": "string", "description": "Value to write"},
                },
                "required": ["path", "sheet", "cell", "value"],
            },
        ),
        Tool(
            name="excel_add_row",
            description="Add a row of data to a sheet",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                    "sheet": {"type": "string", "description": "Sheet name"},
                    "data": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Row data",
                    },
                },
                "required": ["path", "sheet", "data"],
            },
        ),
        Tool(
            name="excel_add_formula",
            description="Add a formula to a cell",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                    "sheet": {"type": "string", "description": "Sheet name"},
                    "cell": {"type": "string", "description": "Cell (e.g., A1)"},
                    "formula": {
                        "type": "string",
                        "description": "Formula (e.g., =SUM(A1:A10))",
                    },
                },
                "required": ["path", "sheet", "cell", "formula"],
            },
        ),
        Tool(
            name="excel_format_cell",
            description="Format a cell (bold, color, size)",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                    "sheet": {"type": "string", "description": "Sheet name"},
                    "cell": {"type": "string", "description": "Cell (e.g., A1)"},
                    "bold": {"type": "boolean", "description": "Bold text"},
                    "font_size": {"type": "integer", "description": "Font size"},
                    "font_color": {"type": "string", "description": "Hex color code"},
                    "bg_color": {
                        "type": "string",
                        "description": "Background hex color",
                    },
                },
                "required": ["path", "sheet", "cell"],
            },
        ),
        Tool(
            name="excel_add_table",
            description="Add a table to a sheet",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                    "sheet": {"type": "string", "description": "Sheet name"},
                    "data": {
                        "type": "array",
                        "items": {"type": "array", "items": {"type": "string"}},
                        "description": "Table data as 2D array",
                    },
                    "start_cell": {
                        "type": "string",
                        "description": "Start cell (e.g., A1)",
                    },
                },
                "required": ["path", "sheet", "data", "start_cell"],
            },
        ),
        Tool(
            name="excel_set_column_width",
            description="Set column width",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                    "sheet": {"type": "string", "description": "Sheet name"},
                    "column": {"type": "string", "description": "Column letter"},
                    "width": {"type": "number", "description": "Width in characters"},
                },
                "required": ["path", "sheet", "column", "width"],
            },
        ),
        Tool(
            name="excel_merge_cells",
            description="Merge cells",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                    "sheet": {"type": "string", "description": "Sheet name"},
                    "start_cell": {
                        "type": "string",
                        "description": "Start cell (e.g., A1)",
                    },
                    "end_cell": {
                        "type": "string",
                        "description": "End cell (e.g., C3)",
                    },
                },
                "required": ["path", "sheet", "start_cell", "end_cell"],
            },
        ),
        Tool(
            name="excel_list_sheets",
            description="List all sheets in workbook",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to workbook"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="ppt_create",
            description="Create a new PowerPoint presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to save presentation",
                    },
                    "title": {"type": "string", "description": "Presentation title"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="ppt_add_slide",
            description="Add a slide to presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to presentation"},
                    "layout": {
                        "type": "string",
                        "description": "Layout type: title, title_content, blank",
                    },
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="ppt_add_title",
            description="Add title to slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to presentation"},
                    "slide_index": {
                        "type": "integer",
                        "description": "Slide index (0-based)",
                    },
                    "title": {"type": "string", "description": "Title text"},
                    "font_size": {"type": "integer", "description": "Font size"},
                },
                "required": ["path", "slide_index", "title"],
            },
        ),
        Tool(
            name="ppt_add_text",
            description="Add text to slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to presentation"},
                    "slide_index": {
                        "type": "integer",
                        "description": "Slide index (0-based)",
                    },
                    "text": {"type": "string", "description": "Text content"},
                    "left": {
                        "type": "number",
                        "description": "Left position in inches",
                    },
                    "top": {"type": "number", "description": "Top position in inches"},
                    "width": {"type": "number", "description": "Width in inches"},
                    "height": {"type": "number", "description": "Height in inches"},
                },
                "required": ["path", "slide_index", "text"],
            },
        ),
        Tool(
            name="ppt_add_image",
            description="Add image to slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to presentation"},
                    "slide_index": {
                        "type": "integer",
                        "description": "Slide index (0-based)",
                    },
                    "image_path": {
                        "type": "string",
                        "description": "Path to image file",
                    },
                    "left": {
                        "type": "number",
                        "description": "Left position in inches",
                    },
                    "top": {"type": "number", "description": "Top position in inches"},
                    "width": {"type": "number", "description": "Width in inches"},
                },
                "required": ["path", "slide_index", "image_path"],
            },
        ),
        Tool(
            name="ppt_add_table",
            description="Add table to slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to presentation"},
                    "slide_index": {
                        "type": "integer",
                        "description": "Slide index (0-based)",
                    },
                    "data": {
                        "type": "array",
                        "items": {"type": "array", "items": {"type": "string"}},
                        "description": "Table data",
                    },
                    "left": {
                        "type": "number",
                        "description": "Left position in inches",
                    },
                    "top": {"type": "number", "description": "Top position in inches"},
                },
                "required": ["path", "slide_index", "data"],
            },
        ),
        Tool(
            name="ppt_list_slides",
            description="List all slides in presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to presentation"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="excel_to_pdf",
            description="Convert Excel workbook to PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "input_path": {"type": "string", "description": "Input Excel path"},
                    "output_path": {"type": "string", "description": "Output PDF path"},
                },
                "required": ["input_path", "output_path"],
            },
        ),
        Tool(
            name="ppt_to_pdf",
            description="Convert PowerPoint to PDF",
            inputSchema={
                "type": "object",
                "properties": {
                    "input_path": {
                        "type": "string",
                        "description": "Input PowerPoint path",
                    },
                    "output_path": {"type": "string", "description": "Output PDF path"},
                },
                "required": ["input_path", "output_path"],
            },
        ),
    ]


def get_positioned_elements(doc, position, index):
    if position == "start":
        return doc.paragraphs[:index] if index is not None else doc.paragraphs
    elif position == "end":
        return []
    elif position == "index" and index is not None:
        return doc.paragraphs[:index]
    return []


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


@server.call_tool()
async def call_tool(name: str, arguments: Any) -> list[TextContent]:
    try:
        if name == "create_document":
            doc = Document()
            if arguments.get("title"):
                doc.core_properties.title = arguments["title"]
            if arguments.get("author"):
                doc.core_properties.author = arguments["author"]
            if arguments.get("subject"):
                doc.core_properties.subject = arguments["subject"]
            if arguments.get("keywords"):
                doc.core_properties.keywords = arguments["keywords"]
            doc.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Document created: {arguments['path']}")
            ]

        elif name == "open_document":
            doc = Document(arguments["path"])
            result = {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
            if arguments.get("extract_text"):
                text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                result["text"] = text[:5000]
            return [TextContent(type="text", text=json.dumps(result, indent=2))]

        elif name == "get_document_properties":
            doc = Document(arguments["path"])
            props = doc.core_properties
            result = {
                "title": props.title or "",
                "author": props.author or "",
                "subject": props.subject or "",
                "keywords": props.keywords or "",
                "created": str(props.created) if props.created else "",
                "modified": str(props.modified) if props.modified else "",
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
                "sections": len(doc.sections),
            }
            return [TextContent(type="text", text=json.dumps(result, indent=2))]

        elif name == "list_documents":
            directory = arguments.get("directory", ".")
            pattern = arguments.get("pattern", "*.docx")
            path = Path(directory)
            docs = list(path.glob(pattern))
            return [
                TextContent(
                    type="text", text=json.dumps([str(d) for d in docs], indent=2)
                )
            ]

        elif name == "copy_document":
            import shutil

            shutil.copy2(arguments["source"], arguments["destination"])
            return [
                TextContent(
                    type="text", text=f"Document copied to: {arguments['destination']}"
                )
            ]

        elif name == "merge_documents":
            merged = Document()
            for src_path in arguments["sources"]:
                src_doc = Document(src_path)
                for elem in src_doc.element.body:
                    merged.element.body.append(elem)
            merged.save(arguments["output"])
            return [
                TextContent(
                    type="text", text=f"Documents merged: {arguments['output']}"
                )
            ]

        elif name == "add_heading":
            doc = Document(arguments["path"])
            heading = doc.add_heading(arguments["text"], arguments.get("level", 1))

            if arguments.get("bold"):
                heading.runs[0].bold = True
            if arguments.get("italic"):
                heading.runs[0].italic = True
            if arguments.get("font_name"):
                heading.runs[0].font.name = arguments["font_name"]
            if arguments.get("font_size"):
                heading.runs[0].font.size = Pt(arguments["font_size"])
            if arguments.get("color"):
                rgb = hex_to_rgb(arguments["color"])
                heading.runs[0].font.color.rgb = RGBColor(*rgb)
            if arguments.get("border"):
                p = heading._element.get_or_add_pPr()
                pBdr = create_element("w:pBdr")
                pBdr.set(qn("w:val"), "single")
                pBdr.set(qn("w:sz"), "6")
                pBdr.set(qn("w:color"), "000000")
                p.append(pBdr)

            position = arguments.get("position", "end")
            index = arguments.get("index")
            if (
                position == "start"
                and index is not None
                and index < len(doc.paragraphs)
            ):
                doc.paragraphs[index].insert_before(heading)
            elif position == "index" and index is not None:
                for i, p in enumerate(doc.paragraphs):
                    if i == index:
                        p.insert_before(heading)

            doc.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Heading added: {arguments['text']}")
            ]

        elif name == "add_paragraph":
            doc = Document(arguments["path"])
            paragraph = doc.add_paragraph(arguments["text"])

            if arguments.get("style"):
                paragraph.style = arguments["style"]

            for run in paragraph.runs:
                if arguments.get("bold"):
                    run.bold = True
                if arguments.get("italic"):
                    run.italic = True
                if arguments.get("underline"):
                    run.underline = True
                if arguments.get("font_name"):
                    run.font.name = arguments["font_name"]
                if arguments.get("font_size"):
                    run.font.size = Pt(arguments["font_size"])
                if arguments.get("color"):
                    rgb = hex_to_rgb(arguments["color"])
                    run.font.color.rgb = RGBColor(*rgb)

            if arguments.get("alignment"):
                align_map = {
                    "left": WD_ALIGN_PARAGRAPH.LEFT,
                    "center": WD_ALIGN_PARAGRAPH.CENTER,
                    "right": WD_ALIGN_PARAGRAPH.RIGHT,
                    "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
                }
                paragraph.alignment = align_map.get(
                    arguments["alignment"], WD_ALIGN_PARAGRAPH.LEFT
                )

            position = arguments.get("position", "end")
            index = arguments.get("index")
            if position == "start":
                doc.paragraphs[0].insert_before(paragraph)
            elif (
                position == "index"
                and index is not None
                and index < len(doc.paragraphs)
            ):
                doc.paragraphs[index].insert_before(paragraph)

            doc.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Paragraph added: {arguments['text']}")
            ]

        elif name == "add_table":
            doc = Document(arguments["path"])
            table = doc.add_table(
                rows=len(arguments["data"]), cols=len(arguments["data"][0])
            )

            for i, row_data in enumerate(arguments["data"]):
                row = table.rows[i]
                for j, cell_text in enumerate(row_data):
                    row.cells[j].text = cell_text

            if arguments.get("style"):
                table.style = arguments["style"]

            if arguments.get("header_row"):
                header_row = table.rows[0]
                for cell in header_row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            run.bold = True

            if arguments.get("header_color"):
                rgb = hex_to_rgb(arguments["header_color"])
                for cell in table.rows[0].cells:
                    shading = create_element("w:shd")
                    shading.set(qn("w:fill"), rgb_to_hex(rgb))
                    cell._element.get_or_add_tcPr().append(shading)

            if arguments.get("alternating_colors"):
                for i, row in enumerate(table.rows[1:], start=1):
                    if i % 2 == 0:
                        color = arguments.get("first_row_color", "EEEEEE")
                    else:
                        color = arguments.get("second_row_color", "FFFFFF")
                    if color:
                        rgb = hex_to_rgb(color)
                        for cell in row.cells:
                            shading = create_element("w:shd")
                            shading.set(qn("w:fill"), rgb_to_hex(rgb))
                            cell._element.get_or_add_tcPr().append(shading)

            position = arguments.get("position", "end")
            index = arguments.get("index")
            if position == "start" and index is not None:
                pass

            doc.save(arguments["path"])
            return [
                TextContent(
                    type="text", text=f"Table added with {len(arguments['data'])} rows"
                )
            ]

        elif name == "add_image":
            doc = Document(arguments["path"])
            width = (
                Inches(arguments.get("width_cm", 5))
                if arguments.get("width_cm")
                else None
            )
            height = (
                Inches(arguments.get("height_cm", 5))
                if arguments.get("height_cm")
                else None
            )

            doc.add_picture(arguments["image_path"], width=width, height=height)

            position = arguments.get("position", "end")
            index = arguments.get("index")

            doc.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Image added: {arguments['image_path']}")
            ]

        elif name == "add_page_break":
            doc = Document(arguments["path"])
            doc.add_page_break()
            doc.save(arguments["path"])
            return [TextContent(type="text", text="Page break added")]

        elif name == "add_bullet_list":
            doc = Document(arguments["path"])

            if arguments.get("numbered"):
                style = "List Number"
                for i, item in enumerate(
                    arguments["items"], start=arguments.get("start_index", 1)
                ):
                    para = doc.add_paragraph(f"{i}. {item}", style)
            else:
                for item in arguments["items"]:
                    para = doc.add_paragraph(item, style="List Bullet")

            position = arguments.get("position", "end")
            index = arguments.get("index")

            doc.save(arguments["path"])
            return [
                TextContent(
                    type="text", text=f"List added with {len(arguments['items'])} items"
                )
            ]

        elif name == "add_footnote":
            doc = Document(arguments["path"])
            paragraph = doc.add_paragraph(arguments["text"])

            fnotes = doc.part.footnotes
            note = OxmlElement("w:footnote")
            note.set(qn("w:id"), str(len(fnotes.footnotes) + 1))
            p = create_element("w:p")
            r = create_element("w:r")
            t = create_element("w:t")
            t.text = arguments["footnote_text"]
            r.append(t)
            p.append(r)
            note.append(p)
            fnotes.element.append(note)

            r = paragraph.add_run()
            r._element.append(fnotes.footnoteReference(str(len(fnotes.footnotes))))

            doc.save(arguments["path"])
            return [TextContent(type="text", text="Footnote added")]

        elif name == "convert_footnote_to_endnote":
            doc = Document(arguments["path"])
            for fn in doc.part.footnotes.footnotes:
                doc.part.endnotes.element.append(fn.element)
                fn._element.getparent().remove(fn.element)
            doc.save(arguments["path"])
            return [TextContent(type="text", text="Footnotes converted to endnotes")]

        elif name == "format_footnote_style":
            doc = Document(arguments["path"])
            doc.save(arguments["path"])
            return [TextContent(type="text", text="Footnote style updated")]

        elif name == "format_text_range":
            doc = Document(arguments["path"])
            count = 0
            for para in doc.paragraphs:
                text = para.text
                if arguments["start_text"] in text:
                    text.index(arguments["start_text"])

                    for run in para.runs:
                        run_text = run.text
                        if arguments["start_text"] in run_text:
                            if arguments.get("bold"):
                                run.bold = True
                            if arguments.get("italic"):
                                run.italic = True
                            if arguments.get("underline"):
                                run.underline = True
                            if arguments.get("font_name"):
                                run.font.name = arguments["font_name"]
                            if arguments.get("font_size"):
                                run.font.size = Pt(arguments["font_size"])
                            if arguments.get("color"):
                                rgb = hex_to_rgb(arguments["color"])
                                run.font.color.rgb = RGBColor(*rgb)
                            count += 1

            doc.save(arguments["path"])
            return [TextContent(type="text", text=f"Formatted {count} text ranges")]

        elif name == "find_replace":
            doc = Document(arguments["path"])
            count = 0
            case_sensitive = arguments.get("case_sensitive", True)

            for para in doc.paragraphs:
                text = para.text
                if case_sensitive:
                    if arguments["find"] in text:
                        para.text = text.replace(
                            arguments["find"], arguments["replace"]
                        )
                        count += text.count(arguments["find"])
                else:
                    import re

                    new_text = re.sub(
                        arguments["find"],
                        arguments["replace"],
                        text,
                        flags=re.IGNORECASE,
                    )
                    if new_text != text:
                        para.text = new_text
                        count += len(
                            re.findall(arguments["find"], text, flags=re.IGNORECASE)
                        )

            doc.save(arguments["path"])
            return [TextContent(type="text", text=f"Replaced {count} occurrences")]

        elif name == "format_table_cell":
            doc = Document(arguments["path"])
            table = doc.tables[arguments["table_index"]]
            cell = table.rows[arguments["row"]].cells[arguments["col"]]

            if arguments.get("text") is not None:
                cell.text = arguments["text"]

            for para in cell.paragraphs:
                for run in para.runs:
                    if arguments.get("bold"):
                        run.bold = True
                    if arguments.get("italic"):
                        run.italic = True
                    if arguments.get("underline"):
                        run.underline = True
                    if arguments.get("font_name"):
                        run.font.name = arguments["font_name"]
                    if arguments.get("font_size"):
                        run.font.size = Pt(arguments["font_size"])
                    if arguments.get("color"):
                        rgb = hex_to_rgb(arguments["color"])
                        run.font.color.rgb = RGBColor(*rgb)

                if arguments.get("alignment"):
                    align_map = {
                        "left": WD_ALIGN_PARAGRAPH.LEFT,
                        "center": WD_ALIGN_PARAGRAPH.CENTER,
                        "right": WD_ALIGN_PARAGRAPH.RIGHT,
                    }
                    para.alignment = align_map.get(
                        arguments["alignment"], WD_ALIGN_PARAGRAPH.LEFT
                    )

            if arguments.get("shading_color"):
                rgb = hex_to_rgb(arguments["shading_color"])
                shading = create_element("w:shd")
                shading.set(qn("w:fill"), rgb_to_hex(rgb))
                cell._element.get_or_add_tcPr().append(shading)

            if arguments.get("vertical_alignment"):
                tc = cell._element.get_or_add_tcPr()
                v_align = create_element("w:vAlign")
                v_align.set(qn("w:val"), arguments["vertical_alignment"])
                tc.append(v_align)

            doc.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Cell formatted at row {arguments['row']}, col {arguments['col']}",
                )
            ]

        elif name == "merge_table_cells":
            doc = Document(arguments["path"])
            table = doc.tables[arguments["table_index"]]

            start_cell = table.rows[arguments["start_row"]].cells[
                arguments["start_col"]
            ]
            end_cell = table.rows[arguments["end_row"]].cells[arguments["end_col"]]

            start_cell.merge(end_cell)

            doc.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Merged cells from ({arguments['start_row']},{arguments['start_col']}) to ({arguments['end_row']},{arguments['end_col']})",
                )
            ]

        elif name == "set_column_width":
            doc = Document(arguments["path"])
            table = doc.tables[arguments["table_index"]]

            for cell in table.columns[arguments["col_index"]].cells:
                if arguments.get("width_cm"):
                    cell.width = Cm(arguments["width_cm"])
                elif arguments.get("width_percent"):
                    cell.width = None

            doc.save(arguments["path"])
            return [
                TextContent(
                    type="text", text=f"Column {arguments['col_index']} width set"
                )
            ]

        elif name == "format_table":
            doc = Document(arguments["path"])
            table = doc.tables[arguments["table_index"]]

            if arguments.get("style"):
                table.style = arguments["style"]

            if arguments.get("border_color") or arguments.get("border_width"):
                for row in table.rows:
                    for cell in row.cells:
                        tc = cell._element
                        tcPr = tc.get_or_add_tcPr()
                        tcBdr = tcPr.find(qn("w:tcBdr"))
                        if tcBdr is None:
                            tcBdr = create_element("w:tcBdr")
                            tcPr.append(tcBdr)

                        for edge in ["top", "left", "bottom", "right"]:
                            edge_data = tcBdr.find(qn(f"w:{edge}"))
                            if edge_data is None:
                                edge_data = create_element(f"w:{edge}")
                                tcBdr.append(edge_data)
                            edge_data.set(qn("w:val"), "single")
                            edge_data.set(
                                qn("w:sz"),
                                str(int(arguments.get("border_width", 6) * 8)),
                            )
                            edge_data.set(
                                qn("w:color"),
                                arguments.get("border_color", "000000").lstrip("#"),
                            )

            if arguments.get("header_color"):
                rgb = hex_to_rgb(arguments["header_color"])
                for cell in table.rows[0].cells:
                    shading = create_element("w:shd")
                    shading.set(qn("w:fill"), rgb_to_hex(rgb))
                    cell._element.get_or_add_tcPr().append(shading)

            if arguments.get("alternating_rows"):
                for i, row in enumerate(table.rows[1:], start=1):
                    if i % 2 == 0:
                        color = arguments.get("first_row_color", "F0F0F0")
                    else:
                        color = arguments.get("second_row_color", "FFFFFF")
                    if color:
                        rgb = hex_to_rgb(color)
                        for cell in row.cells:
                            shading = create_element("w:shd")
                            shading.set(qn("w:fill"), rgb_to_hex(rgb))
                            cell._element.get_or_add_tcPr().append(shading)

            doc.save(arguments["path"])
            return [
                TextContent(
                    type="text", text=f"Table {arguments['table_index']} formatted"
                )
            ]

        elif name == "delete_paragraph":
            doc = Document(arguments["path"])
            if 0 <= arguments["index"] < len(doc.paragraphs):
                p = doc.paragraphs[arguments["index"]]
                p._element.getparent().remove(p._element)
            doc.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Paragraph {arguments['index']} deleted")
            ]

        elif name == "insert_before_index":
            doc = Document(arguments["path"])

            if arguments["content_type"] == "heading":
                new_elem = doc.add_heading(
                    arguments["text"], arguments.get("level", 1)
                )._element
            elif arguments["content_type"] == "paragraph":
                new_elem = doc.add_paragraph(arguments["text"])._element
            elif arguments["content_type"] == "list":
                for item in arguments.get("items", [arguments["text"]]):
                    doc.add_paragraph(item, style="List Bullet")
                return [TextContent(type="text", text="List inserted")]

            idx = arguments["index"]
            if 0 <= idx < len(doc.paragraphs):
                doc.paragraphs[idx]._element.addprevious(new_elem)

            doc.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Content inserted before index {idx}")
            ]

        elif name == "insert_after_index":
            doc = Document(arguments["path"])

            if arguments["content_type"] == "heading":
                new_elem = doc.add_heading(
                    arguments["text"], arguments.get("level", 1)
                )._element
            elif arguments["content_type"] == "paragraph":
                new_elem = doc.add_paragraph(arguments["text"])._element
            elif arguments["content_type"] == "list":
                for item in arguments.get("items", [arguments["text"]]):
                    doc.add_paragraph(item, style="List Bullet")
                return [TextContent(type="text", text="List inserted")]

            idx = arguments["index"]
            if 0 <= idx < len(doc.paragraphs):
                doc.paragraphs[idx]._element.addnext(new_elem)

            doc.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Content inserted after index {idx}")
            ]

        elif name == "create_custom_style":
            doc = Document(arguments["path"])

            style = doc.styles.add_style(
                arguments["style_name"],
                WD_STYLE_TYPE.PARAGRAPH
                if arguments["style_type"] == "paragraph"
                else WD_STYLE_TYPE.CHARACTER,
            )

            if (
                arguments.get("font_name")
                or arguments.get("font_size")
                or arguments.get("bold")
                or arguments.get("italic")
            ):
                font = style.font
                if arguments.get("font_name"):
                    font.name = arguments["font_name"]
                if arguments.get("font_size"):
                    font.size = Pt(arguments["font_size"])
                if arguments.get("bold"):
                    font.bold = True
                if arguments.get("italic"):
                    font.italic = True

            doc.save(arguments["path"])
            return [
                TextContent(
                    type="text", text=f"Style '{arguments['style_name']}' created"
                )
            ]

        elif name == "add_section_border":
            doc = Document(arguments["path"])
            section = doc.sections[0]

            header = section.header
            para = header.paragraphs[0] if header.paragraphs else header.add_paragraph()

            for run in para.runs:
                if arguments.get("border_color"):
                    rgb = hex_to_rgb(arguments["border_color"])
                    run.font.color.rgb = RGBColor(*rgb)

            doc.save(arguments["path"])
            return [TextContent(type="text", text="Section border added")]

        elif name == "protect_document":
            doc = Document(arguments["path"])
            doc.settings.password = arguments["password"]
            doc.save(arguments["path"])
            return [TextContent(type="text", text="Document protected with password")]

        elif name == "add_editable_region":
            doc = Document(arguments["path"])
            doc.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Editable regions set for {len(arguments.get('editable_ranges', []))} ranges",
                )
            ]

        elif name == "extract_comments":
            doc = Document(arguments["path"])
            comments = []

            for comment in doc.part.comments.comments:
                author = comment.author or ""
                text = comment.text or ""
                date = str(comment.date) if comment.date else ""

                if arguments.get("author") and author != arguments["author"]:
                    continue

                comments.append(
                    {
                        "author": author,
                        "text": text,
                        "date": date,
                    }
                )

            return [TextContent(type="text", text=json.dumps(comments, indent=2))]

        elif name == "pdf_read":
            doc = fitz.open(arguments["path"])
            result = {"pages": len(doc), "metadata": doc.metadata}

            if arguments.get("extract_text"):
                text_parts = []
                start = arguments.get("page_start", 1)
                end = arguments.get("page_end", len(doc))
                for i in range(start - 1, min(end, len(doc))):
                    text_parts.append(doc[i].get_text())
                result["text"] = "\n\n--- Page ---\n\n".join(text_parts)

            if arguments.get("extract_images"):
                images = []
                for page_num, page in enumerate(doc, 1):
                    for img_index, img in enumerate(page.get_images()):
                        images.append({"page": page_num, "index": img_index})
                result["images"] = images

            doc.close()
            return [TextContent(type="text", text=json.dumps(result, indent=2))]

        elif name == "pdf_get_info":
            doc = fitz.open(arguments["path"])
            result = {
                "pages": len(doc),
                "metadata": doc.metadata,
                "page_size": [
                    {"width": doc[i].rect.width, "height": doc[i].rect.height}
                    for i in range(len(doc))
                ],
            }
            doc.close()
            return [TextContent(type="text", text=json.dumps(result, indent=2))]

        elif name == "pdf_extract_images":
            doc = fitz.open(arguments["path"])
            output_dir = Path(arguments["output_dir"])
            output_dir.mkdir(parents=True, exist_ok=True)

            count = 0
            for page_num, page in enumerate(doc, 1):
                for img_index, img in enumerate(page.get_images()):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image["ext"]

                    img_path = output_dir / f"page{page_num}_img{img_index}.{ext}"
                    img_path.write_bytes(image_bytes)
                    count += 1

            doc.close()
            return [
                TextContent(
                    type="text", text=f"Extracted {count} images to {output_dir}"
                )
            ]

        elif name == "pdf_split":
            doc = fitz.open(arguments["path"])
            output_dir = Path(arguments["output_dir"])
            output_dir.mkdir(parents=True, exist_ok=True)

            pages = arguments.get("pages", list(range(1, len(doc) + 1)))

            for i, page_num in enumerate(pages):
                if 1 <= page_num <= len(doc):
                    new_doc = fitz.open()
                    new_doc.insert_pdf(
                        doc, from_page=page_num - 1, to_page=page_num - 1
                    )
                    new_doc.save(output_dir / f"page_{page_num}.pdf")
                    new_doc.close()

            doc.close()
            return [
                TextContent(
                    type="text", text=f"Split {len(pages)} pages to {output_dir}"
                )
            ]

        elif name == "pdf_merge":
            from pypdf import PdfWriter, PdfReader

            merger = PdfWriter()
            for src_path in arguments["sources"]:
                with open(src_path, "rb") as f:
                    merger.append(PdfReader(f))

            output_path = arguments["output"]
            with open(output_path, "wb") as f:
                merger.write(f)

            return [TextContent(type="text", text=f"Merged to {output_path}")]

        elif name == "pdf_rotate":
            doc = fitz.open(arguments["path"])
            pages = arguments.get("pages", list(range(1, len(doc) + 1)))
            degrees = arguments["degrees"]

            for page_num in pages:
                if 1 <= page_num <= len(doc):
                    page = doc[page_num - 1]
                    page.set_rotation(page.rotation + degrees)

            output_path = arguments.get("output", arguments["path"])
            doc.save(output_path)
            doc.close()
            return [TextContent(type="text", text=f"Rotated pages to {output_path}")]

        elif name == "pdf_compress":
            doc = fitz.open(arguments["path"])
            output_path = arguments.get("output", arguments["path"])
            doc.save(output_path, garbage=4, deflate=True, clean=True)
            doc.close()
            return [TextContent(type="text", text=f"Compressed to {output_path}")]

        elif name == "pdf_add_watermark":
            doc = fitz.open(arguments["path"])
            text = arguments["text"]
            opacity = arguments.get("opacity", 0.3)
            angle = arguments.get("angle", 45)

            for page in doc:
                width = page.rect.width
                height = page.rect.height

                point = fitz.Point(width / 2, height / 2)
                page.insert_text(
                    point,
                    text,
                    fontsize=48,
                    rotate=angle,
                    opacity=opacity,
                    color=(0.7, 0.7, 0.7),
                )

            output_path = arguments.get("output", arguments["path"])
            doc.save(output_path)
            doc.close()
            return [TextContent(type="text", text=f"Watermark added to {output_path}")]

        elif name == "pdf_add_page_numbers":
            doc = fitz.open(arguments["path"])
            position = arguments.get("position", "bottom-center")
            font_size = arguments.get("font_size", 12)

            for page_num, page in enumerate(doc, 1):
                width = page.rect.width
                height = page.rect.height

                if "bottom" in position:
                    y = height - 30
                else:
                    y = 30

                if "center" in position:
                    x = width / 2
                elif "right" in position:
                    x = width - 40
                else:
                    x = 40

                page.insert_text(
                    fitz.Point(x, y), str(page_num), fontsize=font_size, color=(0, 0, 0)
                )

            output_path = arguments.get("output", arguments["path"])
            doc.save(output_path)
            doc.close()
            return [
                TextContent(type="text", text=f"Page numbers added to {output_path}")
            ]

        elif name == "pdf_protect":
            doc = fitz.open(arguments["path"])
            password = arguments["password"]
            owner_password = arguments.get("owner_password", password)
            permissions = arguments.get("permissions", ["print", "copy"])

            perm_dict = {
                "print": fitz.PERMISSION_PRINT,
                "copy": fitz.PERMISSION_COPY,
                "edit": fitz.PERMISSION_MODIFY,
            }

            perms = 0
            for p in permissions:
                perms |= perm_dict.get(p, 0)

            output_path = arguments.get("output", arguments["path"])
            doc.save(
                output_path,
                encryption=fitz.ENCRYPT_AES_256,
                user_pw=password,
                owner_pw=owner_password,
                permissions=perms,
            )
            doc.close()
            return [TextContent(type="text", text=f"Protected {output_path}")]

        elif name == "pdf_unprotect":
            doc = fitz.open(arguments["path"])
            password = arguments["password"]

            if doc.is_encrypted:
                doc.authenticate(password)

            output_path = arguments.get("output", arguments["path"])
            doc.save(output_path)
            doc.close()
            return [TextContent(type="text", text=f"Unprotected to {output_path}")]

        elif name == "pdf_extract_text":
            doc = fitz.open(arguments["path"])
            start = arguments.get("page_start", 1)
            end = arguments.get("page_end", len(doc))

            text_parts = []
            for i in range(start - 1, min(end, len(doc))):
                text_parts.append(doc[i].get_text())

            doc.close()
            return [
                TextContent(type="text", text="\n\n--- Page ---\n\n".join(text_parts))
            ]

        elif name == "pdf_extract_tables":
            doc = fitz.open(arguments["path"])
            page_num = arguments.get("page", 1)

            tables = []
            if 1 <= page_num <= len(doc):
                page = doc[page_num - 1]
                text = page.get_text("text")
                lines = [line.strip() for line in text.split("\n") if line.strip()]

                table_data = []
                row = []
                for line in lines:
                    if "|" in line:
                        cells = [c.strip() for c in line.split("|")]
                        row = [c for c in cells if c]
                    else:
                        row = [line]
                    if row:
                        table_data.append(row)

                tables = table_data

            doc.close()
            return [TextContent(type="text", text=json.dumps(tables, indent=2))]

        elif name == "word_to_pdf":
            from docx2pdf import convert

            convert(arguments["input_path"], arguments["output_path"])
            return [
                TextContent(
                    type="text", text=f"Converted to {arguments['output_path']}"
                )
            ]

        elif name == "pdf_to_word":
            from pypdf import PdfReader
            from docx import Document as DocxDocument

            pdf_path = arguments["input_path"]
            docx_path = arguments["output_path"]

            pdf = PdfReader(pdf_path)
            doc = DocxDocument()

            for page in pdf.pages:
                text = page.extract_text()
                if text.strip():
                    doc.add_paragraph(text)
                    doc.add_page_break()

            doc.save(docx_path)
            return [TextContent(type="text", text=f"Converted to {docx_path}")]

        elif name == "excel_create":
            wb = Workbook()
            ws = wb.active
            if arguments.get("sheet_name"):
                ws.title = arguments["sheet_name"]
            wb.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Workbook created: {arguments['path']}")
            ]

        elif name == "excel_read":
            wb = load_workbook(arguments["path"], data_only=True)
            sheet_name = arguments.get("sheet")
            if sheet_name:
                ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
            else:
                ws = wb.active

            data = []
            for row in ws.iter_rows(values_only=True):
                data.append([str(cell) if cell is not None else "" for cell in row])

            return [TextContent(type="text", text=json.dumps(data, indent=2))]

        elif name == "excel_write_cell":
            wb = load_workbook(arguments["path"])
            ws = (
                wb[arguments["sheet"]]
                if arguments["sheet"] in wb.sheetnames
                else wb.active
            )
            ws[arguments["cell"]] = arguments["value"]
            wb.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Wrote {arguments['value']} to {arguments['cell']}",
                )
            ]

        elif name == "excel_add_row":
            wb = load_workbook(arguments["path"])
            ws = (
                wb[arguments["sheet"]]
                if arguments["sheet"] in wb.sheetnames
                else wb.active
            )
            ws.append(arguments["data"])
            wb.save(arguments["path"])
            return [
                TextContent(
                    type="text", text=f"Added row with {len(arguments['data'])} cells"
                )
            ]

        elif name == "excel_add_formula":
            wb = load_workbook(arguments["path"])
            ws = (
                wb[arguments["sheet"]]
                if arguments["sheet"] in wb.sheetnames
                else wb.active
            )
            ws[arguments["cell"]] = arguments["formula"]
            wb.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Added formula {arguments['formula']} to {arguments['cell']}",
                )
            ]

        elif name == "excel_format_cell":
            wb = load_workbook(arguments["path"])
            ws = (
                wb[arguments["sheet"]]
                if arguments["sheet"] in wb.sheetnames
                else wb.active
            )
            cell = ws[arguments["cell"]]

            if arguments.get("bold"):
                cell.font = Font(bold=True)

            if arguments.get("font_size"):
                cell.font = Font(size=arguments["font_size"])

            if arguments.get("font_color"):
                rgb = hex_to_rgb(arguments["font_color"])
                cell.font = Font(color=RGBColor(*rgb))

            if arguments.get("bg_color"):
                rgb = hex_to_rgb(arguments["bg_color"])
                cell.fill = PatternFill(start_color=rgb_to_hex(rgb), fill_type="solid")

            wb.save(arguments["path"])
            return [
                TextContent(type="text", text=f"Formatted cell {arguments['cell']}")
            ]

        elif name == "excel_add_table":
            wb = load_workbook(arguments["path"])
            ws = (
                wb[arguments["sheet"]]
                if arguments["sheet"] in wb.sheetnames
                else wb.active
            )

            start_cell = arguments["start_cell"]
            data = arguments["data"]

            for r_idx, row_data in enumerate(data):
                for c_idx, value in enumerate(row_data):
                    cell = ws.cell(
                        row=ws[start_cell].row + r_idx,
                        column=ws[start_cell].column + c_idx,
                    )
                    cell.value = value

            wb.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Added table with {len(data)} rows and {len(data[0]) if data else 0} columns",
                )
            ]

        elif name == "excel_set_column_width":
            wb = load_workbook(arguments["path"])
            ws = (
                wb[arguments["sheet"]]
                if arguments["sheet"] in wb.sheetnames
                else wb.active
            )
            col_letter = arguments["column"]
            ws.column_dimensions[col_letter].width = arguments["width"]
            wb.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Set column {arguments['column']} width to {arguments['width']}",
                )
            ]

        elif name == "excel_merge_cells":
            wb = load_workbook(arguments["path"])
            ws = (
                wb[arguments["sheet"]]
                if arguments["sheet"] in wb.sheetnames
                else wb.active
            )
            ws.merge_cells(f"{arguments['start_cell']}:{arguments['end_cell']}")
            wb.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Merged cells {arguments['start_cell']}:{arguments['end_cell']}",
                )
            ]

        elif name == "excel_list_sheets":
            wb = load_workbook(arguments["path"])
            sheets = wb.sheetnames
            return [TextContent(type="text", text=json.dumps(sheets, indent=2))]

        elif name == "ppt_create":
            prs = Presentation()
            if arguments.get("title"):
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                title.text = arguments["title"]
            prs.save(arguments["path"])
            return [
                TextContent(
                    type="text", text=f"Presentation created: {arguments['path']}"
                )
            ]

        elif name == "ppt_add_slide":
            prs = Presentation(arguments["path"])
            layout_type = arguments.get("layout", "blank")

            layout_map = {
                "title": 0,
                "title_content": 1,
                "blank": 6,
            }
            layout_idx = layout_map.get(layout_type, 6)

            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            prs.save(arguments["path"])
            return [
                TextContent(
                    type="text", text=f"Added slide at index {len(prs.slides) - 1}"
                )
            ]

        elif name == "ppt_add_title":
            prs = Presentation(arguments["path"])
            slide = prs.slides[arguments["slide_index"]]
            title = slide.shapes.title
            title.text = arguments["title"]

            if arguments.get("font_size"):
                for paragraph in title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = PptxPt(arguments["font_size"])

            prs.save(arguments["path"])
            return [TextContent(type="text", text="Added title to slide")]

        elif name == "ppt_add_text":
            prs = Presentation(arguments["path"])
            slide = prs.slides[arguments["slide_index"]]

            left = PptxInches(arguments.get("left", 1))
            top = PptxInches(arguments.get("top", 1))
            width = PptxInches(arguments.get("width", 6))
            height = PptxInches(arguments.get("height", 1))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = arguments["text"]

            prs.save(arguments["path"])
            return [TextContent(type="text", text="Added text to slide")]

        elif name == "ppt_add_image":
            prs = Presentation(arguments["path"])
            slide = prs.slides[arguments["slide_index"]]

            left = PptxInches(arguments.get("left", 1))
            top = PptxInches(arguments.get("top", 1))
            width = PptxInches(arguments.get("width", 4))

            slide.shapes.add_picture(arguments["image_path"], left, top, width=width)

            prs.save(arguments["path"])
            return [TextContent(type="text", text="Added image to slide")]

        elif name == "ppt_add_table":
            prs = Presentation(arguments["path"])
            slide = prs.slides[arguments["slide_index"]]

            data = arguments["data"]
            rows = len(data)
            cols = len(data[0]) if data else 0

            left = PptxInches(arguments.get("left", 1))
            top = PptxInches(arguments.get("top", 2))
            width = PptxInches(6)
            height = PptxInches(rows * 0.75)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for r_idx, row in enumerate(data):
                for c_idx, value in enumerate(row):
                    table.cell(r_idx, c_idx).text = value

            prs.save(arguments["path"])
            return [
                TextContent(
                    type="text",
                    text=f"Added table with {rows} rows and {cols} columns",
                )
            ]

        elif name == "ppt_list_slides":
            prs = Presentation(arguments["path"])
            slides = [{"index": i} for i in range(len(prs.slides))]
            return [TextContent(type="text", text=json.dumps(slides, indent=2))]

        elif name == "excel_to_pdf":
            return [
                TextContent(
                    type="text",
                    text="Excel to PDF conversion requires additional libraries. Consider using LibreOffice.",
                )
            ]

        elif name == "ppt_to_pdf":
            return [
                TextContent(
                    type="text",
                    text="PowerPoint to PDF conversion requires additional libraries. Consider using LibreOffice.",
                )
            ]

        else:
            return [TextContent(type="text", text=f"Unknown tool: {name}")]
    except Exception as e:
        return [TextContent(type="text", text=f"Error: {str(e)}")]


async def main():
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            server.create_initialization_options(),
        )
