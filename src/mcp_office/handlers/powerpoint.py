from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

from mcp_office.handlers.base import DocumentHandler
from mcp.types import Tool, TextContent


class PptHandler(DocumentHandler):
    """Handler for PowerPoint operations"""

    def get_tools(self) -> list[Tool]:
        return [
            Tool(
                name="ppt_create",
                description="Create a new PowerPoint presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Save path"},
                        "title": {"type": "string", "description": "Title"},
                    },
                    "required": ["path"],
                },
            ),
            Tool(
                name="ppt_add_slide",
                description="Add a slide",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "layout": {"type": "string", "description": "Layout type"},
                    },
                    "required": ["path"],
                },
            ),
            Tool(
                name="ppt_add_title",
                description="Add title to slide",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "title": {"type": "string", "description": "Title text"},
                        "font_size": {"type": "integer", "description": "Font size"},
                    },
                    "required": ["path", "slide_index", "title"],
                },
            ),
            Tool(
                name="ppt_add_text",
                description="Add text to slide",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "text": {"type": "string", "description": "Text"},
                        "left": {"type": "number", "description": "Left position"},
                        "top": {"type": "number", "description": "Top position"},
                        "width": {"type": "number", "description": "Width"},
                        "height": {"type": "number", "description": "Height"},
                    },
                    "required": ["path", "slide_index", "text"],
                },
            ),
            Tool(
                name="ppt_add_image",
                description="Add image to slide",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "image_path": {"type": "string", "description": "Image path"},
                        "left": {"type": "number", "description": "Left"},
                        "top": {"type": "number", "description": "Top"},
                        "width": {"type": "number", "description": "Width"},
                    },
                    "required": ["path", "slide_index", "image_path"],
                },
            ),
            Tool(
                name="ppt_add_table",
                description="Add table to slide",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "data": {
                            "type": "array",
                            "items": {"type": "array", "items": {"type": "string"}},
                        },
                        "left": {"type": "number", "description": "Left"},
                        "top": {"type": "number", "description": "Top"},
                    },
                    "required": ["path", "slide_index", "data"],
                },
            ),
            Tool(
                name="ppt_list_slides",
                description="List slides",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                    },
                    "required": ["path"],
                },
            ),
            Tool(
                name="ppt_set_background",
                description="Set slide background color",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "color": {
                            "type": "string",
                            "description": "Hex color (e.g., #FF0000)",
                        },
                    },
                    "required": ["path", "slide_index", "color"],
                },
            ),
            Tool(
                name="ppt_add_shape",
                description="Add shape to slide",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "shape_type": {"type": "string", "description": "Shape type"},
                        "left": {"type": "number", "description": "Left position"},
                        "top": {"type": "number", "description": "Top position"},
                        "width": {"type": "number", "description": "Width"},
                        "height": {"type": "number", "description": "Height"},
                        "color": {"type": "string", "description": "Fill color hex"},
                    },
                    "required": ["path", "slide_index", "shape_type"],
                },
            ),
            Tool(
                name="ppt_format_text",
                description="Format text with colors and styles",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "text": {"type": "string", "description": "Text to format"},
                        "bold": {"type": "boolean", "description": "Bold"},
                        "italic": {"type": "boolean", "description": "Italic"},
                        "font_size": {"type": "integer", "description": "Font size"},
                        "font_color": {
                            "type": "string",
                            "description": "Font color hex",
                        },
                        "bg_color": {
                            "type": "string",
                            "description": "Background color hex",
                        },
                    },
                    "required": ["path", "slide_index", "text"],
                },
            ),
            Tool(
                name="ppt_set_theme",
                description="Apply theme colors to presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "primary_color": {
                            "type": "string",
                            "description": "Primary color hex",
                        },
                    },
                    "required": ["path"],
                },
            ),
            Tool(
                name="ppt_add_chart",
                description="Add chart to slide",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "chart_type": {
                            "type": "string",
                            "description": "Chart type: column, bar, line, pie, area",
                        },
                        "categories": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Categories",
                        },
                        "values": {
                            "type": "array",
                            "items": {"type": "number"},
                            "description": "Values",
                        },
                        "title": {"type": "string", "description": "Chart title"},
                        "left": {"type": "number", "description": "Left position"},
                        "top": {"type": "number", "description": "Top position"},
                        "width": {"type": "number", "description": "Width"},
                        "height": {"type": "number", "description": "Height"},
                    },
                    "required": [
                        "path",
                        "slide_index",
                        "chart_type",
                        "categories",
                        "values",
                    ],
                },
            ),
            Tool(
                name="ppt_add_bullet_slide",
                description="Add bullet slide with formatted text",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "title": {"type": "string", "description": "Slide title"},
                        "bullets": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Bullet points",
                        },
                    },
                    "required": ["path", "title", "bullets"],
                },
            ),
            Tool(
                name="ppt_set_slide_layout",
                description="Set slide layout",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "layout": {
                            "type": "string",
                            "description": "Layout: title, title_content, bullet, blank, title_only",
                        },
                    },
                    "required": ["path", "slide_index", "layout"],
                },
            ),
            Tool(
                name="ppt_add_textbox",
                description="Add formatted textbox",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                        "slide_index": {
                            "type": "integer",
                            "description": "Slide index",
                        },
                        "text": {"type": "string", "description": "Text content"},
                        "left": {"type": "number", "description": "Left position"},
                        "top": {"type": "number", "description": "Top position"},
                        "width": {"type": "number", "description": "Width"},
                        "height": {"type": "number", "description": "Height"},
                        "font_size": {"type": "integer", "description": "Font size"},
                        "bold": {"type": "boolean", "description": "Bold text"},
                        "color": {"type": "string", "description": "Text color hex"},
                    },
                    "required": ["path", "slide_index", "text"],
                },
            ),
            Tool(
                name="ppt_extract_text",
                description="Extract all text from presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Presentation path"},
                    },
                    "required": ["path"],
                },
            ),
        ]

    async def execute(self, tool_name: str, arguments: dict) -> list[TextContent]:
        try:
            if tool_name == "ppt_create":
                return self._create(arguments)
            elif tool_name == "ppt_add_slide":
                return self._add_slide(arguments)
            elif tool_name == "ppt_add_title":
                return self._add_title(arguments)
            elif tool_name == "ppt_add_text":
                return self._add_text(arguments)
            elif tool_name == "ppt_add_image":
                return self._add_image(arguments)
            elif tool_name == "ppt_add_table":
                return self._add_table(arguments)
            elif tool_name == "ppt_list_slides":
                return self._list_slides(arguments)
            elif tool_name == "ppt_set_background":
                return self._set_background(arguments)
            elif tool_name == "ppt_add_shape":
                return self._add_shape(arguments)
            elif tool_name == "ppt_format_text":
                return self._format_text(arguments)
            elif tool_name == "ppt_set_theme":
                return self._set_theme(arguments)
            elif tool_name == "ppt_add_chart":
                return self._add_chart(arguments)
            elif tool_name == "ppt_add_bullet_slide":
                return self._add_bullet_slide(arguments)
            elif tool_name == "ppt_set_slide_layout":
                return self._set_slide_layout(arguments)
            elif tool_name == "ppt_add_textbox":
                return self._add_textbox(arguments)
            elif tool_name == "ppt_extract_text":
                return self._extract_text(arguments)
            return self.error_result(f"Unknown tool: {tool_name}")
        except Exception as e:
            return self.error_result(str(e))

    def _create(self, args: dict) -> list[TextContent]:
        prs = Presentation()
        if title := args.get("title"):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
        prs.save(args["path"])
        return self.success_result(f"Presentation created: {args['path']}")

    def _add_slide(self, args: dict) -> list[TextContent]:
        prs = Presentation(args["path"])
        layout_map = {"title": 0, "title_content": 1, "blank": 6}
        idx = layout_map.get(args.get("layout", "blank"), 6)
        prs.slides.add_slide(prs.slide_layouts[idx])
        prs.save(args["path"])
        return self.success_result(f"Added slide at index {len(prs.slides) - 1}")

    def _add_title(self, args: dict) -> list[TextContent]:
        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]
        title = slide.shapes.title
        title.text = args["title"]
        if fs := args.get("font_size"):
            for p in title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptxPt(fs)
        prs.save(args["path"])
        return self.success_result("Added title to slide")

    def _add_text(self, args: dict) -> list[TextContent]:
        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]
        left = PptxInches(args.get("left", 1))
        top = PptxInches(args.get("top", 1))
        width = PptxInches(args.get("width", 6))
        height = PptxInches(args.get("height", 1))
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text = args["text"]
        prs.save(args["path"])
        return self.success_result("Added text to slide")

    def _add_image(self, args: dict) -> list[TextContent]:
        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]
        left = PptxInches(args.get("left", 1))
        top = PptxInches(args.get("top", 1))
        width = PptxInches(args.get("width", 4))
        slide.shapes.add_picture(args["image_path"], left, top, width=width)
        prs.save(args["path"])
        return self.success_result("Added image to slide")

    def _add_table(self, args: dict) -> list[TextContent]:
        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]
        data = args["data"]
        rows, cols = len(data), len(data[0])
        left = PptxInches(args.get("left", 1))
        top = PptxInches(args.get("top", 2))
        table = slide.shapes.add_table(rows, cols, left, top).table
        for r_idx, row in enumerate(data):
            for c_idx, val in enumerate(row):
                table.cell(r_idx, c_idx).text = val
        prs.save(args["path"])
        return self.success_result(f"Added table {rows}x{cols}")

    def _list_slides(self, args: dict) -> list[TextContent]:
        import json

        prs = Presentation(args["path"])
        slides = [{"index": i} for i in range(len(prs.slides))]
        return [TextContent(type="text", text=json.dumps(slides, indent=2))]

    def _set_background(self, args: dict) -> list[TextContent]:
        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]

        background = slide.background
        fill = background.fill
        fill.solid()

        color = args.get("color", "#FF0000").lstrip("#")
        r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
        fill.fore_color.rgb = PptxRGBColor(r, g, b)

        prs.save(args["path"])
        return self.success_result("Background color set")

    def _add_shape(self, args: dict) -> list[TextContent]:
        from mcp_office.base import hex_to_rgb

        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]

        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "diamond": MSO_SHAPE.DIAMOND,
            "triangle": MSO_SHAPE.TRIANGLE,
            "star": MSO_SHAPE.STAR_5,
            "arrow": MSO_SHAPE.ARROW_RIGHT,
        }

        shape_type = shape_map.get(
            args.get("shape_type", "rectangle").lower(), MSO_SHAPE.RECTANGLE
        )

        left = PptxInches(args.get("left", 1))
        top = PptxInches(args.get("top", 1))
        width = PptxInches(args.get("width", 2))
        height = PptxInches(args.get("height", 2))

        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        if color := args.get("color"):
            rgb = hex_to_rgb(color)
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptxRGBColor(*rgb)

        prs.save(args["path"])
        return self.success_result(f"Shape added: {args.get('shape_type')}")

    def _format_text(self, args: dict) -> list[TextContent]:
        from mcp_office.base import hex_to_rgb

        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                if args["text"] in shape.text:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if args.get("bold"):
                                run.font.bold = True
                            if args.get("italic"):
                                run.font.italic = True
                            if fs := args.get("font_size"):
                                run.font.size = PptxPt(fs)
                            if fc := args.get("font_color"):
                                rgb = hex_to_rgb(fc)
                                run.font.color.rgb = PptxRGBColor(*rgb)

        prs.save(args["path"])
        return self.success_result("Text formatted")

    def _set_theme(self, args: dict) -> list[TextContent]:
        from mcp_office.base import hex_to_rgb

        prs = Presentation(args["path"])

        primary = args.get("primary_color", "#FF0000")
        primary_rgb = hex_to_rgb(primary)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = PptxRGBColor(*primary_rgb)

        prs.save(args["path"])
        return self.success_result("Theme colors applied")

    def _add_chart(self, args: dict) -> list[TextContent]:

        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]

        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
        }

        chart_type = chart_type_map.get(
            args.get("chart_type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED
        )

        chart_data = CategoryChartData()
        chart_data.categories = args.get("categories", [])
        chart_data.add_series("Series", args.get("values", []))

        left = PptxInches(args.get("left", 2))
        top = PptxInches(args.get("top", 2))
        width = PptxInches(args.get("width", 6))
        height = PptxInches(args.get("height", 4.5))

        chart = slide.shapes.add_chart(
            chart_type, left, top, width, height, chart_data
        ).chart

        if title := args.get("title"):
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = title

        prs.save(args["path"])
        return self.success_result(f"Chart added: {args.get('chart_type')}")

    def _add_bullet_slide(self, args: dict) -> list[TextContent]:
        prs = Presentation(args["path"])
        bullet_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_layout)

        slide.shapes.title.text = args["title"]

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = args["bullets"][0]

        for bullet in args["bullets"][1:]:
            p = tf.add_paragraph()
            p.text = bullet

        prs.save(args["path"])
        return self.success_result(
            f"Bullet slide added with {len(args['bullets'])} points"
        )

    def _set_slide_layout(self, args: dict) -> list[TextContent]:
        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]

        layout_map = {
            "title": 0,
            "title_content": 1,
            "bullet": 1,
            "blank": 6,
            "title_only": 5,
        }

        layout_idx = layout_map.get(args.get("layout", "blank"), 6)
        slide.layout = prs.slide_layouts[layout_idx]

        prs.save(args["path"])
        return self.success_result("Slide layout updated")

    def _add_textbox(self, args: dict) -> list[TextContent]:
        from mcp_office.base import hex_to_rgb

        prs = Presentation(args["path"])
        slide = prs.slides[args["slide_index"]]

        left = PptxInches(args.get("left", 1))
        top = PptxInches(args.get("top", 1))
        width = PptxInches(args.get("width", 6))
        height = PptxInches(args.get("height", 1))

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = args["text"]

        if args.get("bold"):
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.bold = True

        if fs := args.get("font_size"):
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = PptxPt(fs)

        if color := args.get("color"):
            rgb = hex_to_rgb(color)
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = PptxRGBColor(*rgb)

        prs.save(args["path"])
        return self.success_result("Textbox added")

    def _extract_text(self, args: dict) -> list[TextContent]:
        import json

        prs = Presentation(args["path"])
        all_text = []

        for i, slide in enumerate(prs.slides):
            slide_text = {"slide": i, "text": []}
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text["text"].append(para.text)
            all_text.append(slide_text)

        return [TextContent(type="text", text=json.dumps(all_text, indent=2))]
